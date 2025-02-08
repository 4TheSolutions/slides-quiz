from pptx import Presentation
from pptx.util import Inches
import openai

openai.api_key = "YOUR API KEY"
# models = openai.Model.list()
# print([model["id"] for model in models["data"]])



def generate_quiz_question(bullet_point):
    prompt = f"Generate a short-answer quiz question based solely on the following key point from a biology lecture. Ensure the question is specific to the given bullet point and can be answered using only the provided information. Do not include an answer: '{bullet_point}'."

    response = openai.ChatCompletion.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}]
    )

    return response["choices"][0]["message"]["content"].strip()

# Example
'''
example = "Mitochondria are the powerhouse of the cell."
quiz_question = generate_quiz_question(example)
print(quiz_question)
'''

presentation = Presentation("bioslides.pptx")

slides_content = []

for slide in presentation.slides:
    slide_data = {"bullet_points": []}
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    slide_data["bullet_points"].append(text)
    slides_content.append(slide_data)

'''
for i, slide in enumerate(slides_content):
    print(f"Slide {i+1}:")
    for bp in slide["bullet_points"]:
        print(f"- {bp}")
        '''

from pptx import Presentation
from pptx.util import Inches, Pt

# Load the existing PowerPoint file
ppt_path = "bioslides.pptx"
presentation = Presentation(ppt_path)

# Function to insert a quiz slide before the lecture slide
def insert_quiz_slide(prs, slide_index, bullet_points):
    """Inserts a quiz slide BEFORE a given lecture slide with wrapped text."""
    # Create a new quiz slide
    quiz_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout

    # Add the quiz title
    title = quiz_slide.shapes.title
    if title:
        title.text = f"Quiz Slide {slide_index + 1}"  # Set the title

    # Create a textbox for questions
    left, top, width, height = Inches(1), Inches(1.5), Inches(8.5), Inches(5)
    quiz_textbox = quiz_slide.shapes.add_textbox(left, top, width, height)
    quiz_text_frame = quiz_textbox.text_frame
    quiz_text_frame.word_wrap = True  # Enable text wrapping

    # Generate quiz questions and add them to the textbox
    for bullet_point in bullet_points:
        quiz_question = generate_quiz_question(bullet_point)  # Generate a question from the bullet
        paragraph = quiz_text_frame.add_paragraph()
        paragraph.text = f"• {quiz_question}"
        paragraph.space_after = Pt(10)  # Add spacing between questions

    # Move the newly created quiz slide to the correct position (before the lecture slide)
    xml_slides = prs.slides._sldIdLst  # Access the slide XML structure
    slides = list(xml_slides)
    quiz_slide_element = slides[-1]  # Get the last added slide (our quiz slide)
    
    # Move quiz slide before the lecture slide
    xml_slides.remove(quiz_slide_element)
    xml_slides.insert(slide_index, quiz_slide_element)

# Process each slide and insert quiz slides BEFORE the lecture slides
for i in range(len(presentation.slides)):  
    slide = presentation.slides[i]

    # Extract bullet points from the slide
    bullet_points = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    bullet_points.append(text)

    # Insert the quiz slide before this lecture slide
    insert_quiz_slide(presentation, i, bullet_points)

# Save the updated PowerPoint presentation
presentation.save("bioslides_updated.pptx")  # Save as a new file to preserve the original
